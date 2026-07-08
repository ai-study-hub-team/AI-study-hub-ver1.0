import logging
import os
import time
from google import genai
from settings import GEMINI_API_KEY, GEMINI_MODEL
import random
from google.genai import errors
from google.genai import types
import concurrent.futures
logger = logging.getLogger("ai-service.text_extractor")

IMAGE_PROMPT = """Bạn hãy phân tích bức ảnh này và trả về kết quả bằng tiếng Việt theo cấu trúc sau:
- Tổng quan nội dung ảnh
- Văn bản xuất hiện trong ảnh nếu có
- Mô tả chi tiết ảnh/sơ đồ/slide/bài học nếu có
- Ý nghĩa học tập của ảnh
- English keywords hỗ trợ semantic search

Nếu ảnh không rõ hoặc không có chữ, không được bịa. Hãy mô tả đúng những gì thấy được."""

VIDEO_PROMPT = """Bạn hãy phân tích video này và trả về kết quả bằng tiếng Việt theo cấu trúc sau:
- Tóm tắt video
- Nội dung theo mốc thời gian nếu Gemini xác định được
- Văn bản/slide xuất hiện trong video nếu có
- Nội dung âm thanh/giọng nói nếu nghe được
- Mô tả cảnh/hành động/đối tượng
- English keywords hỗ trợ semantic search

Nếu âm thanh hoặc hình ảnh không rõ thì ghi rõ [không nghe rõ] hoặc [không nhìn rõ], không được bịa."""

AUDIO_PROMPT = """Bạn hãy phân tích file âm thanh này và trả về kết quả bằng tiếng Việt theo cấu trúc sau:
- Tóm tắt nội dung audio
- Transcript hoặc nội dung nghe được
- Các ý chính
- English keywords hỗ trợ semantic search

Nếu không nghe rõ thì ghi rõ [không nghe rõ]."""

def get_gemini_client():
    if not GEMINI_API_KEY:
        raise ValueError("GEMINI_API_KEY is not set")
    return genai.Client(api_key=GEMINI_API_KEY)

def guess_media_mime_type(file_path: str, file_type: str) -> str:
    file_type = file_type.lower().strip()
    if '/' in file_type and (file_type.startswith('image/') or file_type.startswith('video/') or file_type.startswith('audio/')):
        return file_type
        
    _, ext = os.path.splitext(file_path)
    ext = ext.lower().strip('.')
    
    mapping = {
        'png': 'image/png',
        'jpg': 'image/jpeg',
        'jpeg': 'image/jpeg',
        'webp': 'image/webp',
        'gif': 'image/gif',
        'mp4': 'video/mp4',
        'mov': 'video/quicktime',
        'avi': 'video/x-msvideo',
        'mkv': 'video/x-matroska',
        'mp3': 'audio/mp3',
        'wav': 'audio/wav',
        'm4a': 'audio/x-m4a',
        'ogg': 'audio/ogg'
    }
    
    if ext in mapping:
        return mapping[ext]
        
    import mimetypes
    guessed, _ = mimetypes.guess_type(file_path)
    if guessed and (guessed.startswith('image/') or guessed.startswith('video/') or guessed.startswith('audio/')):
        return guessed
        
    if ext in ['png', 'jpg', 'jpeg', 'webp', 'gif']:
        return 'image/jpeg'
    elif ext in ['mp4', 'mov', 'avi', 'mkv']:
        return 'video/mp4'
    elif ext in ['mp3', 'wav', 'm4a', 'ogg']:
        return 'audio/mp3'
        
    return file_type

def wait_for_gemini_file_active(
    client,
    uploaded_file,
    timeout_seconds: int = 180,
    poll_interval_seconds: int = 2,
    max_get_retries: int = 8
):
    """
    Wait until Gemini uploaded file becomes ACTIVE.

    This version handles temporary Gemini File API errors:
    - 500 INTERNAL
    - 503 UNAVAILABLE
    - 429 RESOURCE_EXHAUSTED

    These errors can happen while Gemini is still processing the uploaded video/audio.
    """

    file_name = uploaded_file.name
    start_time = time.time()
    get_error_count = 0

    logger.info(f"Waiting for Gemini file resource {file_name} to be ACTIVE...")

    while time.time() - start_time < timeout_seconds:
        try:
            file_info = client.files.get(name=file_name)
            get_error_count = 0

        except errors.ServerError as e:
            error_text = str(e)
            get_error_count += 1

            if (
                "500" in error_text
                or "503" in error_text
                or "INTERNAL" in error_text
                or "UNAVAILABLE" in error_text
            ):
                if get_error_count > max_get_retries:
                    raise RuntimeError(
                        f"Gemini file status check failed too many times. "
                        f"file={file_name}, error={repr(e)}"
                    )

                wait_seconds = min(3 * get_error_count, 20) + random.randint(0, 3)

                logger.warning(
                    f"Temporary Gemini File API error while checking file status. "
                    f"Retry {get_error_count}/{max_get_retries} after {wait_seconds}s. "
                    f"Error: {repr(e)}"
                )

                time.sleep(wait_seconds)
                continue

            raise

        except errors.APIError as e:
            error_text = str(e)
            get_error_count += 1

            if (
                "429" in error_text
                or "RESOURCE_EXHAUSTED" in error_text
                or "500" in error_text
                or "503" in error_text
            ):
                if get_error_count > max_get_retries:
                    raise RuntimeError(
                        f"Gemini file status check failed too many times. "
                        f"file={file_name}, error={repr(e)}"
                    )

                wait_seconds = min(5 * get_error_count, 30) + random.randint(0, 5)

                logger.warning(
                    f"Temporary Gemini API error while checking file status. "
                    f"Retry {get_error_count}/{max_get_retries} after {wait_seconds}s. "
                    f"Error: {repr(e)}"
                )

                time.sleep(wait_seconds)
                continue

            raise

        state = getattr(file_info, "state", None)
        state_name = getattr(state, "name", str(state))

        logger.info(f"Gemini file {file_name} current state: {state_name}")

        if "ACTIVE" in state_name:
            logger.info(f"Gemini file resource {file_name} is ACTIVE.")
            return file_info

        if "FAILED" in state_name:
            raise RuntimeError(f"Gemini file processing failed. file={file_name}")

        logger.info(
            f"Gemini file is processing, waiting {poll_interval_seconds} seconds..."
        )

        time.sleep(poll_interval_seconds)

    raise TimeoutError(
        f"Timed out waiting for Gemini file {file_name} to become ACTIVE "
        f"after {timeout_seconds} seconds."
    )

def extract_media_with_gemini(file_path: str, mime_type: str) -> str:
    client = get_gemini_client()
    uploaded_file = None
    try:
        logger.info(f"Uploading {file_path} to Gemini (mime_type={mime_type})...")
        uploaded_file = client.files.upload(file=file_path)
        
        # Poll state until active
        uploaded_file = wait_for_gemini_file_active(client, uploaded_file)
        
        if mime_type.startswith("video/"):
            prompt = VIDEO_PROMPT
        elif mime_type.startswith("audio/"):
            prompt = AUDIO_PROMPT
        else:
            raise ValueError(f"Unsupported Gemini media mime type: {mime_type}")
            
        logger.info(f"Sending analysis request to model: {GEMINI_MODEL}")
        response = client.models.generate_content(
            model=GEMINI_MODEL,
            contents=[uploaded_file, prompt]
        )
        
        text = response.text
        if not text or not text.strip():
            raise ValueError("Gemini returned empty text")
            
        return text
    finally:
        if uploaded_file:
            try:
                logger.info(f"Cleaning up Gemini uploaded file: {uploaded_file.name}")
                client.files.delete(name=uploaded_file.name)
            except Exception as delete_err:
                logger.warning(f"Failed to delete uploaded Gemini file resource {uploaded_file.name}: {delete_err}")

def extract_image_inline(file_path: str, mime_type: str, prompt: str = IMAGE_PROMPT) -> str:
    client = get_gemini_client()
    try:
        logger.info(f"Sending inline image analysis request to Gemini (mime_type={mime_type})...")
        with open(file_path, "rb") as f:
            image_bytes = f.read()
            
        image_part = types.Part.from_bytes(
            data=image_bytes,
            mime_type=mime_type,
        )
        
        response = client.models.generate_content(
            model=GEMINI_MODEL,
            contents=[image_part, prompt]
        )
        
        text = response.text
        if not text or not text.strip():
            raise ValueError("Gemini returned empty text")
            
        return text
    except Exception as e:
        logger.error(f"Failed to process inline image with Gemini: {e}")
        raise

def extract_text(file_path: str, file_type: str) -> str:
    """Extracts text from a local file based on its type."""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
        
    # Standardize the file type format
    file_type = file_type.lower().strip()
    if file_type.startswith('.'):
        file_type = file_type[1:]
        
    logger.info(f"Extracting text from {file_path} of type {file_type}")
    
    _, ext = os.path.splitext(file_path)
    ext = ext.lower().strip('.')
    
    # Detect if it's image, video, audio by mime prefix or file extension
    is_image = file_type.startswith('image/') or ext in ['png', 'jpg', 'jpeg', 'webp', 'gif']
    is_video = file_type.startswith('video/') or ext in ['mp4', 'mov', 'avi', 'mkv']
    is_audio = file_type.startswith('audio/') or ext in ['mp3', 'wav', 'm4a', 'ogg']
    
    if is_image:
        mime_type = guess_media_mime_type(file_path, file_type)
        return extract_image_inline(file_path, mime_type)
    elif is_video or is_audio:
        mime_type = guess_media_mime_type(file_path, file_type)
        return extract_media_with_gemini(file_path, mime_type)
        
    # Try generic mime types handling as well
    if 'pdf' in file_type or ext == 'pdf':
        return extract_text_from_pdf(file_path)
    elif 'excel' in file_type or 'spreadsheet' in file_type or ext in ['xls', 'xlsx']:
        return extract_text_from_excel(file_path, file_type)
    elif 'text' in file_type or 'txt' in file_type or ext == 'txt':
        return extract_text_from_txt(file_path)

    elif 'powerpoint' in file_type or 'presentation' in file_type or ext in ['ppt', 'pptx']:
        return extract_text_from_pptx(file_path)
    elif 'document' in file_type or 'docx' in file_type or ext == 'docx':
        return extract_text_from_docx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type} / {ext}")

def extract_text_from_txt(file_path: str) -> str:
    # Try different encodings just in case
    for encoding in ['utf-8', 'latin-1', 'cp1252']:
        try:
            with open(file_path, 'r', encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    
    # If all fail, open with replace error handler
    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
        return f.read()

def extract_text_from_pdf(file_path: str) -> dict:
    """
    Extract text from a PDF preserving per-page character offsets.

    Returns a dict:
        {
            "text": str,          # full concatenated text (all pages)
            "page_char_map": [    # one entry per page (0-indexed list, 1-based page_number)
                {
                    "page_number": int,   # 1-based
                    "char_start": int,    # inclusive start offset in 'text'
                    "char_end": int,      # exclusive end offset in 'text'
                }
            ]
        }

    The caller (main.py /process-document) uses page_char_map to assign
    locatorType="PAGE", locatorStart, locatorEnd to each chunk.
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise ImportError("PyMuPDF is required for PDF extraction. Please install it using 'pip install PyMuPDF'")

    full_text = ""
    page_char_map = []
    with fitz.open(file_path) as doc:
        for page_num_0based, page in enumerate(doc):
            page_text = page.get_text()
            char_start = len(full_text)
            full_text += page_text
            char_end = len(full_text)
            page_char_map.append({
                "page_number": page_num_0based + 1,  # 1-based
                "char_start": char_start,
                "char_end": char_end,
            })

    return {"text": full_text, "page_char_map": page_char_map}


def extract_text_from_docx(file_path: str) -> str:
    try:
        import docx
    except ImportError:
        raise ImportError(
            "python-docx is required for DOCX extraction. "
            "Please install it using 'pip install python-docx'"
        )

    doc = docx.Document(file_path)

    text_parts = []

    # Paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text.strip())

    # Tables
    for table in doc.tables:
        for row in table.rows:
            row_text = []

            for cell in row.cells:
                cell_text = cell.text.strip()

                if cell_text:
                    row_text.append(cell_text)

            if row_text:
                text_parts.append(" | ".join(row_text))

    extracted_text = "\n".join(text_parts)
    logger.info(
        f"DOCX extracted {len(extracted_text)} characters "
        f"from {len(doc.paragraphs)} paragraphs "
        f"and {len(doc.tables)} tables"
    )
    return extracted_text


def extract_text_from_excel(file_path: str, file_type: str = "") -> str:
    ext = os.path.splitext(file_path)[1].lower().strip('.')
    
    if ext == 'xls':
        try:
            import xlrd
        except ImportError:
            raise ImportError("xlrd is required for .xls files. pip install xlrd")
        
        wb = xlrd.open_workbook(file_path)
        text_parts = []
        
        for sheet in wb.sheets():
            text_parts.append(f"<<<SHEET:{sheet.name}>>>")
            is_header = True
            for rowx in range(sheet.nrows):
                row_values = sheet.row_values(rowx)
                if all(cell == "" for cell in row_values):
                    continue
                row_str_parts = [str(cell).strip().replace("\n", " ") for cell in row_values]
                row_text = " | ".join(row_str_parts)
                
                if is_header:
                    text_parts.append("<<<HEADER>>>")
                    text_parts.append(row_text)
                    is_header = False
                else:
                    text_parts.append(f"<<<ROW:{rowx+1}>>>")
                    text_parts.append(row_text)
                    
        return "\n".join(text_parts)
    
    else:
        try:
            import openpyxl
        except ImportError:
            raise ImportError("openpyxl is required. pip install openpyxl")
        
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"<<<SHEET:{sheet_name}>>>")
            is_header = True
            for idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                if all(cell is None or str(cell).strip() == "" for cell in row):
                    continue
                row_str_parts = [str(cell).strip().replace("\n", " ") if cell is not None else "" for cell in row]
                row_text = " | ".join(row_str_parts)
                if is_header:
                    text_parts.append("<<<HEADER>>>")
                    text_parts.append(row_text)
                    is_header = False
                else:
                    text_parts.append(f"<<<ROW:{idx}>>>")
                    text_parts.append(row_text)
                    
        return "\n".join(text_parts)


PPTX_IMAGE_PROMPT = """Bạn là Kiến trúc sư Hệ thống và Chuyên gia Phân tích. Hãy trả về kết quả theo format sau:
DÒNG 1: [TYPE] -> Chọn 1 trong 4: DIAGRAM, CHART, PHOTO, LOGO.
DÒNG 2 TRỞ ĐI: [ANALYSIS]
- Nếu là DIAGRAM (Sơ đồ): Phân tích Purpose, Main Actors, Input, Output, Components, Relationships, Data Flow, Key Concepts. Không mô tả màu sắc/hình dạng.
- Nếu là CHART (Biểu đồ): Mô tả Loại biểu đồ, Trục X, Trục Y, Legend, Peak, Lowest Point, Trend, Trend Change, Anomaly.
- Nếu là PHOTO (Ảnh/Screenshot): Mô tả tổng quan nội dung, văn bản trong ảnh, ý nghĩa học tập.
- Nếu là LOGO (Logo/Icon/Watermark): Trả về mô tả ngắn gọn."""

def extract_text_from_pptx(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower().strip('.')
    if ext == 'ppt':
        raise ValueError("Định dạng .ppt cũ không được hỗ trợ để bóc tách. Vui lòng chuyển đổi sang .pptx trước khi upload.")
        
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        raise ImportError("python-pptx is required. pip install python-pptx")
        
    try:
        import imagehash
        from PIL import Image
        import io
        import hashlib
        import tempfile
    except ImportError:
        raise ImportError("Pillow and ImageHash are required. pip install Pillow ImageHash")

    prs = Presentation(file_path)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Phase 1: Scan Entire Presentation
    parsed_slides = []
    image_records = []
    phash_counter = {}
    
    for idx, slide in enumerate(prs.slides, start=1):
        slide_types = set()
        slide_parts = []
        
        slide_title = ""
        if slide.shapes.title:
            slide_title = slide.shapes.title.text.strip()
            
        if slide_title:
            slide_parts.append(f"Title: {slide_title}")
            
        picture_count_in_slide = 0
        
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                text = shape.text.strip()
                if text and text != slide_title:
                    slide_parts.append(text)
                    if "public class " in text or "def " in text or "SELECT " in text.upper():
                        slide_types.add("CODE")
                    else:
                        slide_types.add("TEXT")
                        
            elif shape.has_table:
                slide_types.add("TABLE")
                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip().replace("\n", " "))
                    slide_parts.append(" | ".join(row_data))
                    
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_count_in_slide += 1
                is_first_picture = (picture_count_in_slide == 1)
                
                image_blob = shape.image.blob
                img = Image.open(io.BytesIO(image_blob))
                h = str(imagehash.phash(img))
                
                phash_counter[h] = phash_counter.get(h, 0) + 1
                
                width_px = shape.width.inches * 96 if shape.width else 0
                height_px = shape.height.inches * 96 if shape.height else 0
                
                width_ratio = shape.width / slide_width if slide_width and shape.width else 0
                height_ratio = shape.height / slide_height if slide_height and shape.height else 0
                
                img_ext = shape.image.ext
                
                record = {
                    "slide_index": idx,
                    "shape_index": shape_idx,
                    "is_first_picture": is_first_picture,
                    "width_px": width_px,
                    "height_px": height_px,
                    "width_ratio": width_ratio,
                    "height_ratio": height_ratio,
                    "image_size": len(image_blob),
                    "phash": h,
                    "image_blob": image_blob,
                    "ext": img_ext
                }
                image_records.append(record)
                
                slide_parts.append({
                    "job_type": "image",
                    "record": record
                })
                
            elif shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                media_blob = None
                media_ext = "mp4"
                try:
                    xml = shape._element.xml
                    import re
                    match = re.search(r'<(?:a:videoFile|a:audioFile)[^>]+r:link="(rId\d+)"', xml)
                    if match:
                        r_id = match.group(1)
                        if r_id in shape.part.rels:
                            rel = shape.part.rels[r_id]
                            if rel.target_part:
                                media_blob = rel.target_part.blob
                                if hasattr(rel.target_part, 'partname') and rel.target_part.partname:
                                    media_ext = rel.target_part.partname.ext.strip('.')
                except Exception as e:
                    logger.warning(f"Failed to extract media blob: {e}")
                    
                if media_blob:
                    h = hashlib.md5(media_blob).hexdigest()
                    import tempfile
                    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{media_ext}") as tmp:
                        tmp.write(media_blob)
                        tmp_path = tmp.name
                        
                    is_audio = media_ext.lower() in ['mp3', 'wav', 'm4a', 'ogg', 'wma']
                    job_type = "audio" if is_audio else "video"
                    
                    slide_parts.append({
                        "job_type": job_type,
                        "hash": h,
                        "tmp_path": tmp_path,
                        "ext": media_ext
                    })
                else:
                    slide_parts.append("[Embedded Media: Lỗi trích xuất]")
                    slide_types.add("MIXED")
                
        parsed_slides.append({
            "idx": idx,
            "parts": slide_parts,
            "slide_types": slide_types
        })

    # Phase 2: Classification
    stats = {
        "total": len(image_records),
        "small": 0,
        "background": 0,
        "template": 0,
        "cache": 0,
        "gemini": 0
    }

    media_cache = {} 
    unique_image_jobs = []
    unique_media_jobs = []
    
    for slide_data in parsed_slides:
        for part in slide_data["parts"]:
            if isinstance(part, dict) and "job_type" in part:
                if part["job_type"] == "image":
                    record = part["record"]
                    h = record["phash"]
                    idx = record["slide_index"]
                    
                    # Step 1: Small Image Filter
                    if record["width_px"] < 120 or record["height_px"] < 120:
                        logger.info(f"[SKIP][SMALL_IMAGE]\nSlide {idx}")
                        part["classification"] = "SMALL_IMAGE"
                        stats["small"] += 1
                        continue
                        
                    # Step 2: Repeated Image Detection & Step 3: Background vs Template
                    if phash_counter.get(h, 0) >= 3:
                        is_bg_size = record["width_ratio"] > 0.8 and record["height_ratio"] > 0.8
                        is_first = record["is_first_picture"]
                        
                        if is_bg_size and is_first:
                            logger.info(f"[SKIP][BACKGROUND]\nSlide {idx}")
                            part["classification"] = "BACKGROUND"
                            stats["background"] += 1
                        else:
                            logger.info(f"[SKIP][TEMPLATE_GRAPHIC]\nSlide {idx}")
                            part["classification"] = "TEMPLATE_GRAPHIC"
                            stats["template"] += 1
                        continue
                        
                    # Step 4: Duplicate Cache
                    if h in media_cache:
                        logger.info(f"[CACHE]\nReused Gemini result for Slide {idx}")
                        part["classification"] = "CACHE"
                        stats["cache"] += 1
                        continue
                        
                    # Step 5: Gemini Analysis
                    logger.info(f"[GEMINI]\nAnalyzing educational image on Slide {idx}...")
                    part["classification"] = "GEMINI"
                    stats["gemini"] += 1
                    
                    media_cache[h] = None 
                    unique_image_jobs.append(part)
                    
                else: 
                    h = part["hash"]
                    if h in media_cache:
                        part["classification"] = "CACHE"
                    else:
                        part["classification"] = "GEMINI"
                        media_cache[h] = None
                        unique_media_jobs.append(part)

    # Phase 3: Process Media Concurrently
    def process_job(job):
        job_type = job["job_type"]
        tmp_path = None
        
        try:
            if job_type == "image":
                record = job["record"]
                ext = record["ext"]
                h = record["phash"]
                image_blob = record["image_blob"]
                
                with tempfile.NamedTemporaryFile(delete=False, suffix=f".{ext}") as tmp:
                    tmp.write(image_blob)
                    tmp_path = tmp.name
                    
                mime_type = guess_media_mime_type(tmp_path, f"image/{ext}")
                return h, extract_image_inline(tmp_path, mime_type, prompt=PPTX_IMAGE_PROMPT).strip()
                
            elif job_type == "video":
                tmp_path = job["tmp_path"]
                ext = job["ext"]
                h = job["hash"]
                mime_type = guess_media_mime_type(tmp_path, f"video/{ext}")
                return h, extract_media_with_gemini(tmp_path, mime_type).strip()
                
            elif job_type == "audio":
                tmp_path = job["tmp_path"]
                ext = job["ext"]
                h = job["hash"]
                mime_type = guess_media_mime_type(tmp_path, f"audio/{ext}")
                return h, extract_media_with_gemini(tmp_path, mime_type).strip()
                
        except Exception as e:
            logger.error(f"Failed to process {job_type} with Gemini: {e}")
            h_val = job.get("record", {}).get("phash") if job_type == "image" else job.get("hash")
            return h_val, f"[{job_type.capitalize()} extraction failed]"
        finally:
            if tmp_path and os.path.exists(tmp_path):
                try:
                    os.remove(tmp_path)
                except OSError:
                    pass

    all_jobs = unique_image_jobs + unique_media_jobs
    if all_jobs:
        with concurrent.futures.ThreadPoolExecutor(max_workers=5) as executor:
            futures = [executor.submit(process_job, job) for job in all_jobs]
            for future in concurrent.futures.as_completed(futures):
                try:
                    h, text_res = future.result()
                    media_cache[h] = text_res
                except Exception as exc:
                    logger.error(f"Job generated an exception: {exc}")

    # Phase 4: Cache Gemini Result and Log Statistics
    stats_msg = f"""
========== PPT Image Statistics ==========
Total images           : {stats['total']}
Small skipped          : {stats['small']}
Background skipped     : {stats['background']}
Template skipped       : {stats['template']}
Gemini cache reused    : {stats['cache']}
Gemini analyzed        : {stats['gemini']}
========================================="""
    logger.info(stats_msg)

    # Phase 5: Merge results back into the original slide order
    text_parts = []
    
    for slide_data in parsed_slides:
        idx = slide_data["idx"]
        slide_types = slide_data["slide_types"]
        slide_text_parts = []
        
        for part in slide_data["parts"]:
            if isinstance(part, dict) and "job_type" in part:
                if part["job_type"] == "image":
                    classification = part.get("classification")
                    
                    if classification in ["SMALL_IMAGE", "BACKGROUND", "TEMPLATE_GRAPHIC"]:
                        continue # Skip these images
                        
                    h = part["record"]["phash"]
                    text_res = media_cache.get(h)
                    if not text_res:
                        text_res = "[Image extraction failed]"
                        
                    lines = text_res.split("\n", 1)
                    if lines:
                        type_line = lines[0].upper()
                        if "DIAGRAM" in type_line:
                            slide_types.add("DIAGRAM")
                        elif "CHART" in type_line:
                            slide_types.add("CHART")
                        else:
                            slide_types.add("IMAGE")
                            
                    slide_text_parts.append(text_res)
                    
                else:
                    h = part["hash"]
                    text_res = media_cache.get(h)
                    if not text_res:
                        text_res = "[Media extraction failed]"
                    slide_types.add("MIXED")
                    slide_text_parts.append(text_res)
                    
            else:
                slide_text_parts.append(part)
                
        if len(slide_types) == 0:
            slide_type = "TEXT"
        elif len(slide_types) == 1:
            slide_type = list(slide_types)[0]
        else:
            if "DIAGRAM" in slide_types:
                slide_type = "DIAGRAM"
            elif "CHART" in slide_types:
                slide_type = "CHART"
            else:
                slide_type = "MIXED"
                
        text_parts.append(f"<<<SLIDE_START:{idx}|TYPE:{slide_type}>>>")
        text_parts.append(f"[Slide {idx}]")
        text_parts.append("\n\n".join(slide_text_parts))
        
    return "\n".join(text_parts)
